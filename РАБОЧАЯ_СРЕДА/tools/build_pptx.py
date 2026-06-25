"""
build_pptx.py — простой спек-Markdown → .pptx (фирменный стиль проекта).

Формат спека (слайды разделяются строкой `---`):

    # Заголовок титульного слайда
    ## подзаголовок / авторы            (только на первом слайде → титул)

    ---
    # Заголовок обычного слайда
    - тезис
    - тезис
        - вложенный тезис               (отступ = под-пункт)

    ---
    # Слайд с картинкой
    ![](../материалы/charts/site_plan_maly_taldykol.png)
    - подпись / тезис справа

Палитра: teal (#0E7C7B) + red-accent (#D7263D), шрифт Calibri/Arial.

Примеры:
    python tools/build_pptx.py slides/концепт.md
    python tools/build_pptx.py slides/концепт.md -o out/Концепт.pptx
"""
from __future__ import annotations
import argparse
import re
import sys
from pathlib import Path

import common

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

TEAL = RGBColor(0x0E, 0x7C, 0x7B)
RED = RGBColor(0xD7, 0x26, 0x3D)
DARK = RGBColor(0x22, 0x2B, 0x2E)
GREY = RGBColor(0x5A, 0x66, 0x6B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Calibri"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)  # 16:9


def parse_spec(text: str) -> list[dict]:
    """Разбирает спек-markdown в список слайдов."""
    blocks = re.split(r"^\s*---\s*$", text, flags=re.MULTILINE)
    slides = []
    for bi, block in enumerate(blocks):
        lines = [ln for ln in block.splitlines()]
        slide = {"title": "", "subtitle": "", "image": None, "bullets": [], "is_title": bi == 0}
        for ln in lines:
            s = ln.strip()
            if not s:
                continue
            m_img = re.match(r"!\[[^\]]*\]\(([^)]+)\)", s)
            if s.startswith("# "):
                slide["title"] = s[2:].strip()
            elif s.startswith("## "):
                slide["subtitle"] = s[3:].strip()
            elif m_img:
                slide["image"] = m_img.group(1).strip()
            elif s.startswith(("- ", "* ")):
                indent = len(ln) - len(ln.lstrip())
                level = 1 if indent >= 2 else 0
                slide["bullets"].append((level, s[2:].strip()))
            else:
                slide["bullets"].append((0, s))
        if slide["title"] or slide["bullets"] or slide["image"]:
            slides.append(slide)
    return slides


def _add_bg_band(slide):
    """Тонкая верхняя плашка teal для фирменного вида."""
    band = slide.shapes.add_shape(1, 0, 0, EMU_W, Inches(0.22))
    band.fill.solid(); band.fill.fore_color.rgb = TEAL
    band.line.fill.background()
    band.shadow.inherit = False


def render_title(slide, spec, src: Path):
    _add_bg_band(slide)
    box = slide.shapes.add_textbox(Inches(0.9), Inches(2.4), Inches(11.5), Inches(2.6))
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = spec["title"]
    p.font.size = Pt(40); p.font.bold = True; p.font.color.rgb = TEAL; p.font.name = FONT
    if spec["subtitle"]:
        p2 = tf.add_paragraph(); p2.text = spec["subtitle"]
        p2.font.size = Pt(20); p2.font.color.rgb = GREY; p2.font.name = FONT
    bar = slide.shapes.add_shape(1, Inches(0.95), Inches(2.25), Inches(2.2), Inches(0.06))
    bar.fill.solid(); bar.fill.fore_color.rgb = RED; bar.line.fill.background()


def render_content(slide, spec):
    _add_bg_band(slide)
    # заголовок
    head = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(12), Inches(1.0))
    hp = head.text_frame.paragraphs[0]; hp.text = spec["title"]
    hp.font.size = Pt(28); hp.font.bold = True; hp.font.color.rgb = TEAL; hp.font.name = FONT
    underline = slide.shapes.add_shape(1, Inches(0.75), Inches(1.35), Inches(1.6), Inches(0.05))
    underline.fill.solid(); underline.fill.fore_color.rgb = RED; underline.line.fill.background()

    has_img = bool(spec["image"])
    text_w = Inches(6.0) if has_img else Inches(11.9)
    if spec["bullets"]:
        body = slide.shapes.add_textbox(Inches(0.75), Inches(1.7), text_w, Inches(5.2))
        tf = body.text_frame; tf.word_wrap = True
        first = True
        for level, txt in spec["bullets"]:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.text = ("• " if level == 0 else "– ") + txt
            p.level = level
            p.font.size = Pt(18 if level == 0 else 15)
            p.font.color.rgb = DARK if level == 0 else GREY
            p.font.name = FONT
            p.space_after = Pt(6)
    if has_img:
        img_path = (spec["_src_dir"] / spec["image"]).resolve()
        if img_path.exists():
            left = Inches(7.0)
            slide.shapes.add_picture(str(img_path), left, Inches(1.7),
                                     width=Inches(5.6))
        else:
            note = slide.shapes.add_textbox(Inches(7.0), Inches(3.2), Inches(5.6), Inches(1))
            note.text_frame.paragraphs[0].text = f"[нет файла: {spec['image']}]"
            note.text_frame.paragraphs[0].font.color.rgb = RED


def build(src: Path, dst: Path | None) -> Path:
    text = src.read_text(encoding="utf-8")
    specs = parse_spec(text)
    if not specs:
        raise SystemExit("В спеке нет слайдов.")
    prs = Presentation()
    prs.slide_width = EMU_W; prs.slide_height = EMU_H
    blank = prs.slide_layouts[6]
    for i, spec in enumerate(specs):
        spec["_src_dir"] = src.parent
        slide = prs.slides.add_slide(blank)
        if i == 0 and spec.get("is_title") and spec.get("subtitle") is not None and not spec["bullets"]:
            render_title(slide, spec, src)
        elif i == 0 and spec.get("is_title") and not spec["bullets"] and not spec["image"]:
            render_title(slide, spec, src)
        else:
            render_content(slide, spec)
    if dst is None:
        dst = common.OUT / (src.stem + ".pptx")
    dst.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(dst))
    print(f"✓ PPTX ({len(specs)} слайд.) → {dst}")
    return dst


def main():
    ap = argparse.ArgumentParser(description="Спек-Markdown → PPTX (фирменный стиль)")
    ap.add_argument("source", help="входной .md спек слайдов")
    ap.add_argument("-o", "--out", help="выходной .pptx (по умолчанию out/<имя>.pptx)")
    args = ap.parse_args()
    build(Path(args.source).resolve(),
          Path(args.out).resolve() if args.out else None)


if __name__ == "__main__":
    sys.exit(main())
